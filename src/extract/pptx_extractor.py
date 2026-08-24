from __future__ import annotations

import os

os.environ["TORCH_COMPILE_DISABLE"] = "1"
os.environ["PYTORCH_JIT"] = "0"
os.environ["OMP_NUM_THREADS"] = "1"

import importlib.util
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.extract.base import (
    BaseExtractor,
    DocumentSection,
    ExtractedDocument,
    ExtractedImage,
    generate_document_id,
)

logger = logging.getLogger(__name__)
_HAS_DOCLING = importlib.util.find_spec("docling") is not None


class PPTXExtractor(BaseExtractor):
    def extract(self, file_path: Path, role: str, category: str) -> ExtractedDocument:
        if _HAS_DOCLING:
            docling_result = self._try_docling_extract(file_path, role, category)
            if docling_result:
                return docling_result

        return self._pptx_extract(file_path, role, category)


    @staticmethod
    def _reading_order_sort_key(shape):
        """
        Key hàm sắp xếp shapes theo thứ tự đọc tự nhiên:
        top → bottom (hàng), cùng hàng thì left → right.

        Đơn vị EMU: 1 inch = 914,400 EMU.
        Bucket 200,000 EMU ≈ 5.5mm — đủ rộng để gom shapes cùng dòng
        mà không lẫn lộn các dòng khác nhau.
        """
        top = getattr(shape, "top", None) or 0
        left = getattr(shape, "left", None) or 0
        row_bucket = top // 200_000  # Gom shapes cùng "hàng" vào 1 bucket
        return (row_bucket, left)

    def _try_docling_extract(self, file_path: Path, role: str, category: str) -> ExtractedDocument | None:
        try:
            from docling.document_converter import DocumentConverter

            converter = DocumentConverter()
            res = converter.convert(str(file_path))
            docling_doc = res.document
            markdown_content = docling_doc.export_to_markdown()

            doc_id = generate_document_id(category, file_path.name)
            title = file_path.stem.replace("_", " ").replace("-", " ")
            num_slides = len(docling_doc.pages) if hasattr(docling_doc, "pages") and docling_doc.pages else 1

            sections = [
                DocumentSection(
                    title=title,
                    level=1,
                    content=markdown_content,
                    section_type="slide",
                )
            ]

            return ExtractedDocument(
                document_id=doc_id,
                title=title,
                source_file=file_path.name,
                source_path=f"{category}/{file_path.name}",
                document_type="pptx",
                role=role,
                category=category,
                pages=num_slides,
                slides=num_slides,
                sections=sections,
                images=[],
                raw_text=markdown_content,
            )
        except BaseException as e:
            logger.warning(f"Docling extraction skipped/failed for PPTX {file_path.name}: {e}")
            return None

    def _pptx_extract(self, file_path: Path, role: str, category: str) -> ExtractedDocument:
        prs = Presentation(file_path)
        sections: list[DocumentSection] = []
        images: list[ExtractedImage] = []
        full_raw_text: list[str] = []

        total_slides = len(prs.slides)

        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            slide_title = ""
            slide_texts: list[str] = []
            slide_tables: list[str] = []
            img_in_slide_count = 0

            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()

            # Sắp xếp shapes theo thứ tự đọc tự nhiên (top→bottom, left→right)
            # để đảm bảo text box, caption và ảnh nhúng được xử lý đúng thứ tự.
            sorted_shapes = sorted(
                slide.shapes,
                key=self._reading_order_sort_key,
            )

            for shape in sorted_shapes:
                if shape.has_text_frame:
                    if shape != slide.shapes.title:
                        text = shape.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.has_table:
                    table_md = self._convert_pptx_table_to_markdown(shape.table)
                    if table_md:
                        slide_tables.append(table_md)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_in_slide_count += 1
                    image_bytes = shape.image.blob
                    ext = shape.image.ext
                    img_name = f"{file_path.stem}_s{slide_num}_img{img_in_slide_count}.{ext}"
                    images.append(
                        ExtractedImage(
                            filename=img_name,
                            image_bytes=image_bytes,
                            slide_num=slide_num,
                        )
                    )

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            section_content_parts = []
            if slide_title:
                section_content_parts.append(f"### Title\n\n{slide_title}")

            if slide_texts:
                section_content_parts.append("### Content\n\n" + "\n\n".join(slide_texts))

            if slide_tables:
                section_content_parts.append("### Tables\n\n" + "\n\n".join(slide_tables))

            if notes_text:
                section_content_parts.append(f"### Notes\n\n{notes_text}")

            full_section_str = "\n\n".join(section_content_parts)
            full_raw_text.append(full_section_str)

            header_title = f"Slide {slide_num}: {slide_title}" if slide_title else f"Slide {slide_num}"
            sections.append(
                DocumentSection(
                    title=header_title,
                    level=2,
                    content=full_section_str,
                    section_type="slide",
                    slide_num=slide_num,
                )
            )

        doc_id = generate_document_id(category, file_path.name)
        title = file_path.stem.replace("_", " ").replace("-", " ")

        return ExtractedDocument(
            document_id=doc_id,
            title=title,
            source_file=file_path.name,
            source_path=f"{category}/{file_path.name}",
            document_type="pptx",
            role=role,
            category=category,
            pages=total_slides,
            slides=total_slides,
            sections=sections,
            images=images,
            raw_text="\n\n".join(full_raw_text),
        )

    def _convert_pptx_table_to_markdown(self, table) -> str:
        rows_data = []
        for row in table.rows:
            row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows_data.append(row_cells)

        if not rows_data:
            return ""

        headers = rows_data[0]
        markdown_lines = [
            "| " + " | ".join(headers) + " |",
            "| " + " | ".join(["---"] * len(headers)) + " |",
        ]

        for row in rows_data[1:]:
            if len(row) < len(headers):
                row.extend([""] * (len(headers) - len(row)))
            elif len(row) > len(headers):
                row = row[: len(headers)]
            markdown_lines.append("| " + " | ".join(row) + " |")

        return "\n".join(markdown_lines)
