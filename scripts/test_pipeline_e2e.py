import json
from pathlib import Path

import docx
import openpyxl
from pptx import Presentation

from src.preprocess.pipeline import PreprocessingPipeline


def create_sample_files(raw_dir: Path):
    raw_dir.mkdir(parents=True, exist_ok=True)

    # 1. Sale - DOCX with PII
    sale_dir = raw_dir / "Sale"
    sale_dir.mkdir(parents=True, exist_ok=True)
    doc = docx.Document()
    doc.add_heading("Quy Trình Bán Hàng", level=1)
    doc.add_paragraph("Khách hàng Nguyễn Văn A")
    doc.add_paragraph("SĐT: 0912345678")
    doc.add_paragraph("Email: nguyenvana@gmail.com")
    doc.save(sale_dir / "QuyTrinhBanHang.docx")

    # 2. KeToan - XLSX with PII
    ketoan_dir = raw_dir / "KeToan"
    ketoan_dir.mkdir(parents=True, exist_ok=True)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "BaoCaoTaiChinh"
    ws.append(["Mã Khách Hàng", "Tên Khách Hàng", "Số Tài Khoản", "Số Điện Thoại"])
    ws.append(["KH001", "Trần Thị B", "1234567890", "0987654321"])
    wb.save(ketoan_dir / "BaoCaoTaiChinh.xlsx")

    # 3. KTV - PPTX
    ktv_dir = raw_dir / "KTV"
    ktv_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hướng Dẫn Kỹ Thuật"
    slide.placeholders[1].text = "Liên hệ KTV Lê Văn C - SĐT: 0901234567"
    prs.save(ktv_dir / "HuongDanKyThuat.pptx")


def run_e2e():
    root = Path(__file__).resolve().parent.parent
    raw_dir = root / "data" / "sample_raw"
    processed_dir = root / "data" / "sample_processed"

    create_sample_files(raw_dir)

    pipeline = PreprocessingPipeline(raw_dir=raw_dir, processed_dir=processed_dir)
    results = pipeline.run_all()

    print("\n==========================================")
    print(f"Successfully processed {len(results)} sample test files.")
    print("==========================================")
    for md_path, meta_path, pii_path in results:
        print(f"\n[Processed File]: {md_path.stem}")
        print(f" - Markdown:   {md_path}")
        print(f" - Metadata:   {meta_path}")
        print(f" - PII Report: {pii_path}")

        md_content = md_path.read_text(encoding="utf-8")
        meta_content = json.loads(meta_path.read_text(encoding="utf-8"))
        pii_content = json.loads(pii_path.read_text(encoding="utf-8"))

        print(f"   Role:          {meta_content.get('role')}")
        print(f"   Access Scope:  {meta_content.get('access_scope')}")
        print(f"   File Hash:     {meta_content.get('file_hash')[:25]}...")
        print(f"   PII Detected:  {pii_content.get('pii_detected')}")
        print(f"   PII Removed:   {pii_content.get('removed_entities')}")

        # Verify no raw PII in output
        for secret in ["0912345678", "nguyenvana@gmail.com", "0987654321"]:
            assert secret not in md_content, f"RAW PII LEAKED IN MD: {secret}"
            assert secret not in json.dumps(meta_content), f"RAW PII LEAKED IN META: {secret}"
            assert secret not in json.dumps(pii_content), f"RAW PII LEAKED IN PII REPORT: {secret}"

    print("\n==========================================")
    print("E2E VERIFICATION PASSED: Pipeline output compliant with SPEC.md!")
    print("==========================================")


if __name__ == "__main__":
    run_e2e()
